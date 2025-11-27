"""Utilities to extract text from uploaded documents."""

from __future__ import annotations

import csv
import io
import zipfile
from pathlib import Path
from typing import Callable

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

TextExtractor = Callable[[bytes], str]


def _extract_text_from_txt(data: bytes) -> str:
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1", errors="ignore")


def _extract_text_from_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text)
    return "\n\n".join(parts)


def _extract_text_from_pptx(data: bytes) -> str:
    presentation = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    parts.append(text)
    return "\n\n".join(parts)


def _extract_text_from_xlsx(data: bytes) -> str:
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def _extract_text_from_csv(data: bytes) -> str:
    text = _extract_text_from_txt(data)
    reader = csv.reader(io.StringIO(text))
    return "\n".join(", ".join(row) for row in reader)


def _extract_text_from_json(data: bytes) -> str:
    """Extract text from JSON files, preserving structure for metrics analysis."""
    text = _extract_text_from_txt(data)
    try:
        # Try to parse and pretty-print JSON for better readability
        import json
        parsed = json.loads(text)
        return json.dumps(parsed, indent=2)
    except json.JSONDecodeError:
        # If not valid JSON, return as-is
        return text


def _extract_text_from_log(data: bytes) -> str:
    """Extract text from log files, preserving log structure."""
    return _extract_text_from_txt(data)


def _extract_text_from_docx(data: bytes) -> str:
    document = Document(io.BytesIO(data))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n".join(parts)


EXTRACTORS: dict[str, TextExtractor] = {
    ".txt": _extract_text_from_txt,
    ".md": _extract_text_from_txt,
    ".log": _extract_text_from_log,
    ".csv": _extract_text_from_csv,
    ".tsv": _extract_text_from_csv,
    ".json": _extract_text_from_json,
    ".pdf": _extract_text_from_pdf,
    ".pptx": _extract_text_from_pptx,
    ".xlsx": _extract_text_from_xlsx,
    ".docx": _extract_text_from_docx,
}


def detect_file_type(filename: str) -> str:
    """Detect file type for metadata purposes."""
    suffix = Path(filename.lower()).suffix
    if suffix in [".log", ".txt"]:
        # Try to detect if it's a log file by checking common log patterns
        return "log"
    elif suffix == ".json":
        return "json"
    elif suffix in [".csv", ".tsv"]:
        return "csv"
    elif suffix == ".pdf":
        return "pdf"
    elif suffix in [".xlsx", ".xls"]:
        return "excel"
    elif suffix == ".docx":
        return "docx"
    elif suffix == ".pptx":
        return "pptx"
    else:
        return "text"


def extract_text(filename: str, data: bytes) -> str:
    suffix = Path(filename.lower()).suffix
    extractor = EXTRACTORS.get(suffix)
    if not extractor:
        raise ValueError(f"Unsupported file type: {suffix or 'unknown'}")
    text = extractor(data)
    if not text.strip():
        raise ValueError("No extractable text found in document.")
    return text


def get_file_metadata(filename: str) -> dict:
    """Get metadata for a file based on its type."""
    file_type = detect_file_type(filename)
    metadata = {
        "filename": filename,
        "file_type": file_type,
    }
    return metadata

